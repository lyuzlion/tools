import os
import math
from io import BytesIO
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import tempfile

def pdf_to_pptx(pdf_path, pptx_path, dpi=300):
    """
    将PDF转换为PPTX，每页对应一张幻灯片，保持原始比例
    
    参数:
    pdf_path: 输入的PDF文件路径
    pptx_path: 输出的PPTX文件路径
    dpi: 图像分辨率（默认300）
    """
    # 创建临时目录
    with tempfile.TemporaryDirectory() as temp_dir:
        # 将PDF转换为图片列表
        images = convert_from_path(
            pdf_path, 
            dpi=dpi,
            output_folder=temp_dir,
            fmt='png',
            use_cropbox=True,
            thread_count=4
        )
        
        if not images:
            raise ValueError("无法从PDF生成图像，请检查PDF文件路径")
        
        # 计算所有页面的最大尺寸（以像素为单位）
        max_width = max(img.width for img in images)
        max_height = max(img.height for img in images)
        
        # 计算最大宽高比
        max_aspect_ratio = max(img.width / img.height for img in images)
        
        # 转换为英寸（用于PPTX）
        max_width_inch = max_width / dpi
        max_height_inch = max_height / dpi
        
        # 创建新的PPTX演示文稿
        prs = Presentation()
        
        # 设置幻灯片尺寸为最大页面尺寸
        prs.slide_width = Inches(max_width_inch)
        prs.slide_height = Inches(max_height_inch)
        
        for i, img in enumerate(images):
            # 创建新幻灯片（使用空白布局）
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 计算当前页面的宽高
            img_width_inch = img.width / dpi
            img_height_inch = img.height / dpi
            
            # 计算居中位置
            left = (max_width_inch - img_width_inch) / 2
            top = (max_height_inch - img_height_inch) / 2
            
            # 将图像保存到字节流
            img_bytes = BytesIO()
            img.save(img_bytes, format='PNG', dpi=(dpi, dpi))
            img_bytes.seek(0)
            
            # 添加图像到幻灯片（保持原始比例）
            pic = slide.shapes.add_picture(
                img_bytes, 
                Inches(left), 
                Inches(top),
                width=Inches(img_width_inch),
                height=Inches(img_height_inch)
            )
            
            print(f"已处理页面 {i+1}/{len(images)} - 原始尺寸: {img.width}x{img.height} 像素")
        
        # 保存PPTX文件
        prs.save(pptx_path)
        print(f"转换完成! 输出文件: {pptx_path}")
        print(f"PPTX尺寸: {max_width_inch:.2f} × {max_height_inch:.2f} 英寸")
        print(f"最大宽高比: {max_aspect_ratio:.4f}")

if __name__ == "__main__":
    # 使用示例
    pdf_to_pptx(
        pdf_path="刘子龙.pdf",    # 输入PDF路径
        pptx_path="output.pptx", # 输出PPTX路径
        dpi=300                 # 分辨率(推荐300-600)
    )